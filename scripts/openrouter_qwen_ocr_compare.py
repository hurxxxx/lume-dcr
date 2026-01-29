#!/usr/bin/env python3
import argparse
import base64
import json
import os
import time
from collections import Counter
from pathlib import Path
from typing import Iterable

import requests
from PIL import Image

OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def load_env_from_file(env_path: Path) -> bool:
    if not env_path.exists():
        return False
    for line in env_path.read_text(encoding="utf-8", errors="ignore").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = value
    return True


def load_env() -> None:
    candidates = [Path(".env"), Path("..") / ".env", Path("../..") / ".env"]
    for candidate in candidates:
        if load_env_from_file(candidate):
            break


def convert_office_to_pdf(src: Path, out_dir: Path) -> Path:
    import subprocess

    ensure_dir(out_dir)
    cmd = [
        "soffice", "--headless", "--nologo", "--nolockcheck", "--nodefault",
        "--convert-to", "pdf", "--outdir", str(out_dir), str(src),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    pdf_path = out_dir / f"{src.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not created for {src}")
    return pdf_path


def convert_pdf_to_images(pdf_path: Path, out_dir: Path, dpi: int) -> list[Path]:
    import subprocess

    ensure_dir(out_dir)
    prefix = out_dir / "page"
    cmd = ["pdftoppm", "-r", str(dpi), "-png", str(pdf_path), str(prefix)]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return sorted(out_dir.glob("page-*.png"))


def extract_pdf_text(pdf_path: Path) -> str:
    import fitz

    lines = []
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text = page.get_text("text") or ""
            lines.append(text.strip())
    return "\n".join([l for l in lines if l])


def extract_pptx_text(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                lines.append(text)
    return "\n".join(lines)


def normalize(text: str) -> str:
    text = text.replace("\u200b", " ").replace("\ufeff", " ")
    return " ".join(text.split())


def edit_distance(a: list, b: list) -> int:
    if not a:
        return len(b)
    if not b:
        return len(a)
    dp = list(range(len(b) + 1))
    for i, ca in enumerate(a, start=1):
        prev = dp[0]
        dp[0] = i
        for j, cb in enumerate(b, start=1):
            cost = 0 if ca == cb else 1
            temp = dp[j]
            dp[j] = min(dp[j] + 1, dp[j - 1] + 1, prev + cost)
            prev = temp
    return dp[-1]


def cer(ref: str, hyp: str) -> float | None:
    ref = normalize(ref)
    hyp = normalize(hyp)
    if not ref:
        return None
    return edit_distance(list(ref), list(hyp)) / len(ref)


def wer(ref: str, hyp: str) -> float | None:
    ref_words = normalize(ref).split()
    hyp_words = normalize(hyp).split()
    if not ref_words:
        return None
    return edit_distance(ref_words, hyp_words) / len(ref_words)


def line_coverage(ref: str, hyp: str) -> float | None:
    ref_lines = [normalize(l) for l in ref.splitlines() if normalize(l)]
    hyp_norm = normalize(hyp)
    if not ref_lines:
        return None
    matched = 0
    for line in ref_lines:
        if len(line) < 4:
            continue
        if line in hyp_norm:
            matched += 1
    return matched / len(ref_lines)


def numeric_recall(ref: str, hyp: str) -> float | None:
    import re

    number_re = re.compile(r"\d+(?:[.,]\d+)?")
    ref_nums = number_re.findall(ref)
    hyp_nums = number_re.findall(hyp)
    if not ref_nums:
        return None
    ref_counts = Counter(ref_nums)
    hyp_counts = Counter(hyp_nums)
    matched = sum(min(count, hyp_counts.get(num, 0)) for num, count in ref_counts.items())
    return matched / len(ref_nums)


def image_to_data_url(image_path: Path, max_side: int) -> str:
    image = Image.open(image_path).convert("RGB")
    if max_side and max(image.size) > max_side:
        image.thumbnail((max_side, max_side), Image.BICUBIC)
    buffer = Path(image_path).with_suffix(".tmp.jpg")
    image.save(buffer, format="JPEG", quality=92)
    data = buffer.read_bytes()
    buffer.unlink(missing_ok=True)
    b64 = base64.b64encode(data).decode("ascii")
    return f"data:image/jpeg;base64,{b64}"


def call_openrouter(api_key: str, model: str, image_url: str, prompt: str, max_tokens: int, timeout: int, retries: int) -> dict:
    headers = {
        "Authorization": f"Bearer {api_key}",
        "X-Title": "lume-dcr-ocr-compare",
    }
    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": image_url}},
                    {"type": "text", "text": prompt},
                ],
            }
        ],
        "max_tokens": max_tokens,
        "temperature": 0,
    }
    last_err = None
    for attempt in range(retries + 1):
        try:
            resp = requests.post(OPENROUTER_URL, headers=headers, json=payload, timeout=timeout)
            if resp.status_code >= 400:
                raise RuntimeError(f"{resp.status_code}: {resp.text[:300]}")
            return resp.json()
        except Exception as exc:
            last_err = exc
            time.sleep(1 + attempt)
    raise RuntimeError(f"OpenRouter call failed: {last_err}")


def extract_text_from_response(data: dict) -> str:
    choices = data.get("choices") or []
    if not choices:
        return ""
    message = choices[0].get("message") or {}
    content = message.get("content")
    if isinstance(content, str):
        return content
    # Some providers may return list content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                parts.append(item.get("text", ""))
        return "\n".join(parts)
    return ""


def write_document(page_paths: Iterable[Path], out_path: Path) -> None:
    parts = []
    for idx, page_path in enumerate(page_paths, start=1):
        text = page_path.read_text(encoding="utf-8", errors="ignore").strip()
        parts.append(f"## Page {idx}\n\n{text}\n")
    out_path.write_text("\n".join(parts).strip() + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="OpenRouter Qwen OCR comparison")
    parser.add_argument("--input", action="append", required=True, help="Input PDF/PPTX path")
    parser.add_argument("--out", type=str, required=True, help="Output root directory")
    parser.add_argument("--models", action="append", help="OpenRouter model id (repeatable)")
    parser.add_argument("--dpi", type=int, default=200)
    parser.add_argument("--max-pages", type=int, default=0)
    parser.add_argument("--max-tokens", type=int, default=1024)
    parser.add_argument("--max-image-side", type=int, default=1280)
    parser.add_argument("--timeout", type=int, default=120)
    parser.add_argument("--retries", type=int, default=2)
    args = parser.parse_args()

    load_env()
    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        raise SystemExit("OPENROUTER_API_KEY not set in environment or .env")

    models = args.models or [
        "qwen/qwen3-vl-8b-instruct",
        "qwen/qwen3-vl-30b-a3b-instruct",
        "qwen/qwen3-vl-32b-instruct",
        "qwen/qwen-2.5-vl-7b-instruct",
    ]

    out_root = Path(args.out).resolve()
    ensure_dir(out_root)
    tmp_root = out_root / "_tmp"
    ensure_dir(tmp_root)

    summary = {
        "started_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "models": models,
        "documents": [],
    }

    for input_path in [Path(p).resolve() for p in args.input]:
        if not input_path.exists():
            raise SystemExit(f"Input not found: {input_path}")

        work_dir = tmp_root / input_path.stem
        ensure_dir(work_dir)

        if input_path.suffix.lower() == ".pdf":
            pdf_path = input_path
            ref_text = extract_pdf_text(pdf_path)
        elif input_path.suffix.lower() == ".pptx":
            pdf_path = convert_office_to_pdf(input_path, work_dir)
            ref_text = extract_pptx_text(input_path)
        else:
            raise SystemExit(f"Unsupported input: {input_path}")

        images = convert_pdf_to_images(pdf_path, work_dir / "pages", args.dpi)
        if args.max_pages:
            images = images[: args.max_pages]

        doc_entry = {
            "file": str(input_path),
            "pages": [],
            "metrics": {},
        }

        for model in models:
            model_slug = model.replace("/", "_").replace(":", "_")
            page_paths = []
            for idx, image_path in enumerate(images, start=1):
                page_out = out_root / model_slug / input_path.stem / f"page_{idx:03d}"
                ensure_dir(page_out)
                image_url = image_to_data_url(image_path, args.max_image_side)
                t0 = time.perf_counter()
                resp = call_openrouter(
                    api_key=api_key,
                    model=model,
                    image_url=image_url,
                    prompt="Convert the document to markdown. Preserve structure and tables if present.",
                    max_tokens=args.max_tokens,
                    timeout=args.timeout,
                    retries=args.retries,
                )
                elapsed = time.perf_counter() - t0
                text = extract_text_from_response(resp)
                result_path = page_out / "result.md"
                result_path.write_text(text, encoding="utf-8")
                (page_out / "response.json").write_text(json.dumps(resp, ensure_ascii=False, indent=2), encoding="utf-8")
                page_paths.append(result_path)
                doc_entry["pages"].append({
                    "model": model,
                    "page": idx,
                    "image": str(image_path),
                    "time_sec": elapsed,
                })
                time.sleep(0.5)

            doc_out = out_root / model_slug / input_path.stem / "document.md"
            write_document(page_paths, doc_out)

            if ref_text:
                hyp = doc_out.read_text(encoding="utf-8", errors="ignore")
                doc_entry["metrics"].setdefault(model, {})
                doc_entry["metrics"][model] = {
                    "cer": cer(ref_text, hyp),
                    "wer": wer(ref_text, hyp),
                    "line_coverage": line_coverage(ref_text, hyp),
                    "numeric_recall": numeric_recall(ref_text, hyp),
                }

        summary["documents"].append(doc_entry)

    summary["finished_at"] = time.strftime("%Y-%m-%d %H:%M:%S")
    (out_root / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")

    report_lines = [
        "# OpenRouter Qwen OCR 비교",
        "",
        f"- Models: {', '.join(models)}",
        f"- Run timestamp: {summary['finished_at']}",
        "",
    ]
    for doc in summary["documents"]:
        report_lines.append(f"## {Path(doc['file']).name}")
        report_lines.append("")
        if doc["metrics"]:
            report_lines.append("| Model | CER | WER | Line coverage | Numeric recall |")
            report_lines.append("| --- | ---: | ---: | ---: | ---: |")
            for model in models:
                m = doc["metrics"].get(model)
                if m:
                    report_lines.append(
                        f"| {model} | {m['cer']:.4f} | {m['wer']:.4f} | {m['line_coverage']:.4f} | {m['numeric_recall']:.4f} |"
                    )
                else:
                    report_lines.append(f"| {model} | n/a | n/a | n/a | n/a |")
            report_lines.append("")
        else:
            report_lines.append("Metrics skipped (missing reference text).")
            report_lines.append("")
        for model in models:
            model_slug = model.replace('/', '_').replace(':', '_')
            report_lines.append(f"- {model}: {out_root / model_slug / Path(doc['file']).stem / 'document.md'}")
        report_lines.append("")

    (out_root / "report.md").write_text("\n".join(report_lines).strip() + "\n", encoding="utf-8")
    print(f"Report written to {out_root / 'report.md'}")


if __name__ == "__main__":
    main()
