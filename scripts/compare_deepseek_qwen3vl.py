#!/usr/bin/env python3
import argparse
import json
import re
import subprocess
import time
from collections import Counter
from pathlib import Path

import torch
from PIL import Image
from transformers import AutoModel, AutoProcessor, AutoTokenizer

DEEPSEEK_MODEL_ID = "deepseek-ai/DeepSeek-OCR-2"
QWEN_MODEL_ID = "Qwen/Qwen3-VL-8B-Instruct"


def run(cmd):
    return subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def convert_office_to_pdf(src: Path, out_dir: Path) -> Path:
    ensure_dir(out_dir)
    cmd = [
        "soffice", "--headless", "--nologo", "--nolockcheck", "--nodefault",
        "--convert-to", "pdf", "--outdir", str(out_dir), str(src),
    ]
    run(cmd)
    pdf_path = out_dir / f"{src.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not created for {src}")
    return pdf_path


def convert_pdf_to_images(pdf_path: Path, out_dir: Path, dpi: int) -> list[Path]:
    ensure_dir(out_dir)
    prefix = out_dir / "page"
    cmd = ["pdftoppm", "-r", str(dpi), "-png", str(pdf_path), str(prefix)]
    run(cmd)
    return sorted(out_dir.glob("page-*.png"))


def extract_pdf_text(pdf_path: Path) -> str:
    import fitz  # PyMuPDF

    lines = []
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text = page.get_text("text") or ""
            lines.append(text.strip())
    return "\n".join([l for l in lines if l])


def extract_pptx_text(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                lines.append(text)
    return "\n".join(lines)


def normalize(text: str) -> str:
    text = text.replace("\u200b", " ").replace("\ufeff", " ")
    return " ".join(text.split())


def edit_distance(a: list, b: list) -> int:
    if not a:
        return len(b)
    if not b:
        return len(a)
    dp = list(range(len(b) + 1))
    for i, ca in enumerate(a, start=1):
        prev = dp[0]
        dp[0] = i
        for j, cb in enumerate(b, start=1):
            cost = 0 if ca == cb else 1
            temp = dp[j]
            dp[j] = min(
                dp[j] + 1,
                dp[j - 1] + 1,
                prev + cost,
            )
            prev = temp
    return dp[-1]


def cer(ref: str, hyp: str) -> float | None:
    ref = normalize(ref)
    hyp = normalize(hyp)
    if not ref:
        return None
    return edit_distance(list(ref), list(hyp)) / len(ref)


def wer(ref: str, hyp: str) -> float | None:
    ref_words = normalize(ref).split()
    hyp_words = normalize(hyp).split()
    if not ref_words:
        return None
    return edit_distance(ref_words, hyp_words) / len(ref_words)


def line_coverage(ref: str, hyp: str) -> float | None:
    ref_lines = [normalize(l) for l in ref.splitlines() if normalize(l)]
    hyp_norm = normalize(hyp)
    if not ref_lines:
        return None
    matched = 0
    for line in ref_lines:
        if len(line) < 4:
            continue
        if line in hyp_norm:
            matched += 1
    return matched / len(ref_lines)


def numeric_recall(ref: str, hyp: str) -> float | None:
    number_re = re.compile(r"\d+(?:[.,]\d+)?")
    ref_nums = number_re.findall(ref)
    hyp_nums = number_re.findall(hyp)
    if not ref_nums:
        return None
    ref_counts = Counter(ref_nums)
    hyp_counts = Counter(hyp_nums)
    matched = sum(min(count, hyp_counts.get(num, 0)) for num, count in ref_counts.items())
    return matched / len(ref_nums)


def load_deepseek():
    dtype = torch.bfloat16 if torch.cuda.is_bf16_supported() else torch.float16
    tokenizer = AutoTokenizer.from_pretrained(DEEPSEEK_MODEL_ID, trust_remote_code=True)
    model = AutoModel.from_pretrained(
        DEEPSEEK_MODEL_ID,
        trust_remote_code=True,
        torch_dtype=dtype,
        device_map="cuda",
        use_cache=True,
        attn_implementation="eager",
    )
    model.eval()
    return tokenizer, model


def load_qwen():
    dtype = torch.bfloat16 if torch.cuda.is_bf16_supported() else torch.float16
    try:
        from transformers import Qwen3VLForConditionalGeneration

        model = Qwen3VLForConditionalGeneration.from_pretrained(
            QWEN_MODEL_ID,
            torch_dtype=dtype,
            device_map="cuda",
        )
    except Exception:
        from transformers import AutoModelForImageTextToText

        model = AutoModelForImageTextToText.from_pretrained(
            QWEN_MODEL_ID,
            torch_dtype=dtype,
            device_map="cuda",
        )
    processor = AutoProcessor.from_pretrained(QWEN_MODEL_ID)
    model.eval()
    return processor, model


def deepseek_infer(tokenizer, model, image_path: Path, out_dir: Path, base_size: int, image_size: int, crop_mode: bool):
    ensure_dir(out_dir)
    prompt = "<image>\n<|grounding|>Convert the document to markdown."
    t0 = time.perf_counter()
    result = model.infer(
        tokenizer,
        prompt=prompt,
        image_file=str(image_path),
        output_path=str(out_dir),
        base_size=base_size,
        image_size=image_size,
        crop_mode=crop_mode,
        save_results=True,
    )
    torch.cuda.synchronize()
    t1 = time.perf_counter()
    result_path = out_dir / "result.mmd"
    return t1 - t0, result_path, result


def qwen_infer(processor, model, image_path: Path, out_dir: Path, max_new_tokens: int, max_image_side: int):
    ensure_dir(out_dir)
    image = Image.open(image_path).convert("RGB")
    if max_image_side and max(image.size) > max_image_side:
        image.thumbnail((max_image_side, max_image_side), Image.BICUBIC)
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "image", "image": image},
                {"type": "text", "text": "Convert the document to markdown."},
            ],
        }
    ]
    prompt = processor.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
    inputs = processor(text=[prompt], images=[image], return_tensors="pt").to(model.device)
    t0 = time.perf_counter()
    generated = model.generate(**inputs, max_new_tokens=max_new_tokens)
    torch.cuda.synchronize()
    t1 = time.perf_counter()
    output = processor.batch_decode(generated, skip_special_tokens=True)[0]
    result_path = out_dir / "result.md"
    result_path.write_text(output, encoding="utf-8")
    return t1 - t0, result_path, output


def write_document(pages: list[Path], out_path: Path) -> None:
    parts = []
    for idx, page_path in enumerate(pages, start=1):
        text = page_path.read_text(encoding="utf-8", errors="ignore").strip()
        parts.append(f"## Page {idx}\n\n{text}\n")
    out_path.write_text("\n".join(parts).strip() + "\n", encoding="utf-8")


def main():
    parser = argparse.ArgumentParser(description="Compare DeepSeek-OCR-2 vs Qwen3-VL OCR outputs.")
    parser.add_argument("--input", action="append", required=True, help="Input PDF/PPTX path")
    parser.add_argument("--out", type=str, required=True, help="Output root directory")
    parser.add_argument("--dpi", type=int, default=200)
    parser.add_argument("--max-pages", type=int, default=0)
    parser.add_argument("--deepseek-base-size", type=int, default=1024)
    parser.add_argument("--deepseek-image-size", type=int, default=768)
    parser.add_argument("--deepseek-crop-mode", action="store_true")
    parser.add_argument("--qwen-max-new-tokens", type=int, default=1024)
    parser.add_argument("--qwen-max-image-side", type=int, default=1280)
    parser.add_argument(
        "--mode",
        choices=["both", "deepseek", "qwen", "report"],
        default="both",
        help="Run inference for deepseek, qwen, both, or just generate report from existing outputs",
    )
    args = parser.parse_args()

    out_root = Path(args.out).resolve()
    ensure_dir(out_root)
    tmp_root = out_root / "_tmp"
    ensure_dir(tmp_root)

    ds_tokenizer = ds_model = None
    qwen_processor = qwen_model = None
    if args.mode in {"both", "deepseek"}:
        print("Loading DeepSeek-OCR-2...")
        ds_tokenizer, ds_model = load_deepseek()
    if args.mode in {"both", "qwen"}:
        print("Loading Qwen3-VL...")
        qwen_processor, qwen_model = load_qwen()

    summary = {
        "started_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "deepseek_model": DEEPSEEK_MODEL_ID,
        "qwen_model": QWEN_MODEL_ID,
        "documents": [],
    }

    for input_path in [Path(p).resolve() for p in args.input]:
        if not input_path.exists():
            raise SystemExit(f"Input not found: {input_path}")

        doc_entry = {
            "file": str(input_path),
            "pages": [],
            "metrics": {},
        }

        work_dir = tmp_root / input_path.stem
        ensure_dir(work_dir)

        if input_path.suffix.lower() == ".pdf":
            pdf_path = input_path
            ref_text = extract_pdf_text(pdf_path)
        elif input_path.suffix.lower() in {".pptx", ".docx"}:
            pdf_path = convert_office_to_pdf(input_path, work_dir)
            ref_text = extract_pptx_text(input_path) if input_path.suffix.lower() == ".pptx" else ""
        else:
            raise SystemExit(f"Unsupported input: {input_path}")

        images = convert_pdf_to_images(pdf_path, work_dir / "pages", args.dpi)
        if args.max_pages:
            images = images[: args.max_pages]

        deepseek_page_paths = []
        qwen_page_paths = []

        for idx, image_path in enumerate(images, start=1):
            page_entry = {"page": idx, "image": str(image_path)}

            if args.mode in {"both", "deepseek"}:
                ds_out = out_root / "deepseek-ocr2" / input_path.stem / f"page_{idx:03d}"
                ds_time, ds_result_path, _ = deepseek_infer(
                    ds_tokenizer,
                    ds_model,
                    image_path,
                    ds_out,
                    args.deepseek_base_size,
                    args.deepseek_image_size,
                    args.deepseek_crop_mode,
                )
                deepseek_page_paths.append(ds_result_path)
                page_entry["deepseek_time_sec"] = ds_time
            else:
                ds_out = out_root / "deepseek-ocr2" / input_path.stem / f"page_{idx:03d}"
                ds_result_path = ds_out / "result.mmd"
                if ds_result_path.exists():
                    deepseek_page_paths.append(ds_result_path)

            if args.mode in {"both", "qwen"}:
                qwen_out = out_root / "qwen3-vl-8b" / input_path.stem / f"page_{idx:03d}"
                qwen_time, qwen_result_path, _ = qwen_infer(
                    qwen_processor,
                    qwen_model,
                    image_path,
                    qwen_out,
                    args.qwen_max_new_tokens,
                    args.qwen_max_image_side,
                )
                qwen_page_paths.append(qwen_result_path)
                page_entry["qwen_time_sec"] = qwen_time
            else:
                qwen_out = out_root / "qwen3-vl-8b" / input_path.stem / f"page_{idx:03d}"
                qwen_result_path = qwen_out / "result.md"
                if qwen_result_path.exists():
                    qwen_page_paths.append(qwen_result_path)

            doc_entry["pages"].append(page_entry)

        deepseek_doc = out_root / "deepseek-ocr2" / input_path.stem / "document.md"
        qwen_doc = out_root / "qwen3-vl-8b" / input_path.stem / "document.md"
        if deepseek_page_paths:
            write_document(deepseek_page_paths, deepseek_doc)
        if qwen_page_paths:
            write_document(qwen_page_paths, qwen_doc)

        metrics = {}
        if ref_text and deepseek_doc.exists() and qwen_doc.exists():
            ds_text = deepseek_doc.read_text(encoding="utf-8", errors="ignore")
            qwen_text = qwen_doc.read_text(encoding="utf-8", errors="ignore")
            metrics = {
                "deepseek": {
                    "cer": cer(ref_text, ds_text),
                    "wer": wer(ref_text, ds_text),
                    "line_coverage": line_coverage(ref_text, ds_text),
                    "numeric_recall": numeric_recall(ref_text, ds_text),
                },
                "qwen3_vl": {
                    "cer": cer(ref_text, qwen_text),
                    "wer": wer(ref_text, qwen_text),
                    "line_coverage": line_coverage(ref_text, qwen_text),
                    "numeric_recall": numeric_recall(ref_text, qwen_text),
                },
            }

        doc_entry["metrics"] = metrics
        summary["documents"].append(doc_entry)

    summary["finished_at"] = time.strftime("%Y-%m-%d %H:%M:%S")
    summary_path = out_root / "summary.json"
    summary_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")

    report_lines = [
        "# DeepSeek-OCR-2 vs Qwen3-VL OCR 비교",
        "",
        f"- DeepSeek model: {DEEPSEEK_MODEL_ID}",
        f"- Qwen model: {QWEN_MODEL_ID}",
        f"- Run timestamp: {summary['finished_at']}",
        "",
    ]
    for doc in summary["documents"]:
        report_lines.append(f"## {Path(doc['file']).name}")
        report_lines.append("")
        if doc["metrics"]:
            report_lines.append("| Model | CER | WER | Line coverage | Numeric recall |")
            report_lines.append("| --- | ---: | ---: | ---: | ---: |")
            for key, label in [("deepseek", "DeepSeek-OCR-2"), ("qwen3_vl", "Qwen3-VL-8B")]:
                if key in doc["metrics"]:
                    m = doc["metrics"][key]
                    report_lines.append(
                        f"| {label} | {m['cer']:.4f} | {m['wer']:.4f} | {m['line_coverage']:.4f} | {m['numeric_recall']:.4f} |"
                    )
                else:
                    report_lines.append(f"| {label} | n/a | n/a | n/a | n/a |")
            report_lines.append("")
        else:
            report_lines.append("Metrics skipped (missing reference text or model outputs).")
            report_lines.append("")
        report_lines.append(
            f"- DeepSeek output: {out_root / 'deepseek-ocr2' / Path(doc['file']).stem / 'document.md'}"
        )
        report_lines.append(
            f"- Qwen3-VL output: {out_root / 'qwen3-vl-8b' / Path(doc['file']).stem / 'document.md'}"
        )
        report_lines.append("")

    (out_root / "report.md").write_text("\n".join(report_lines).strip() + "\n", encoding="utf-8")
    print(f"Report written to {out_root / 'report.md'}")


if __name__ == "__main__":
    main()
